import os
import fitz  # PDF
import pandas as pd
from docx import Document as DocxDocument  
from pptx import Presentation  
from dotenv import load_dotenv
from langchain_community.vectorstores import FAISS
from langchain_openai import AzureOpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document 

load_dotenv()

AZURE_OPENAI_API_KEY = os.environ.get('AZURE_OPENAI_API_KEY')
AZURE_OPENAI_ENDPOINT = os.environ.get("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_VERSION = os.environ.get("AZURE_OPENAI_API_VERSION")
AZURE_OPENAI_EMBEDDING_DEPLOYMENT = os.environ.get("AZURE_OPENAI_EMBEDDING_DEPLOYMENT")


def extract_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    text = ""

    if ext == ".pdf":
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
    
    elif ext == ".docx":
        doc = DocxDocument(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    
    elif ext == ".pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    return text


def build_index_from_file(file_path: str, persist_dir: str = "./rag_faiss_store"):
    full_text = extract_text_from_file(file_path)

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
    documents = text_splitter.split_documents([Document(page_content=full_text)])
    
    embeddings = AzureOpenAIEmbeddings(
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY,
        api_version=AZURE_OPENAI_API_VERSION,
        azure_deployment=AZURE_OPENAI_EMBEDDING_DEPLOYMENT,
    )

    db = FAISS.from_documents(documents, embeddings)
    os.makedirs(persist_dir, exist_ok=True)
    db.save_local(persist_dir)


if __name__ == "__main__":
    # You can change the file below to .docx or .pptx to test other formats
    build_index_from_file("./docs/NehaResumeupdated[1].docx")
